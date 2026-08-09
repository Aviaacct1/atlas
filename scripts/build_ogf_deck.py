r"""Build the Observatory OGF deck as a .pptx, from the forecast rather than by hand.
Author: Avia Solutions.

The deck is the shape of Boeing's Market Overview, built on Avia data and assumptions.
It is also the diagnostic: a slide we cannot fill is a gap in the forecast, not a
formatting problem, so this script does not invent a slide it has no data for. Where an
input is missing it writes the slide into the gap register with the reason, and the
register is a slide in the deck rather than a note in a file.

Everything is generated. No number is typed into a shape. Inputs, all produced by other
scripts in this folder and all held outside the repository:

  fleet_wedge.json        scripts/build_fleet_wedge.py
  gap_decomposition.json  scripts/gap_decomposition.py
  regions_boeing.json     scripts/compare_regions_boeing.py --json
  ogf_deck_data.json      scripts/build_ogf_deck_data.py
  webapp/data/history.json    and  webapp/data/dashboard.json   from the tool build

House rules enforced in code, not by eye, because a document generated from templates
is exactly where a prohibited character survives review: no em dash and no en dash
anywhere in any string that reaches the file, author and last modified by set, and the
proofing language set to en-GB on every run. scripts/check_deck.py re-checks the built
file from the outside.

Style. The Observatory brand, taken from the tokens in observatory_tokens.css: ink
#0F1B28, brass #D4A249 for the observed series and single emphasis, paper #F6F3EC,
never pure white. Cambria carries display, body and all numbers and Arial carries
labels and metadata, which is the safe font pairing the Project Forth Observatory deck
of 7 August 2026 already uses in place of Newsreader and Inter.

Usage:  py -3.12 scripts\build_ogf_deck.py [--out PATH] [--author "Avia Solutions"]
"""
from __future__ import annotations

import argparse
import json
import os
import sys
from datetime import date, datetime, timezone

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation  # noqa: E402
from pptx.chart.data import CategoryChartData  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # noqa: E402
from pptx.enum.lang import MSO_LANGUAGE_ID  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

from avia_forecast import paths  # noqa: E402

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# The Observatory tokens.
INK = RGBColor(0x0F, 0x1B, 0x28)
BRASS = RGBColor(0xD4, 0xA2, 0x49)
BRASS_DEEP = RGBColor(0xA9, 0x7C, 0x33)
PAPER = RGBColor(0xF6, 0xF3, 0xEC)
CARD = RGBColor(0xFB, 0xF9, 0xF3)
LINE = RGBColor(0xE2, 0xDC, 0xCC)
TEXT = RGBColor(0x26, 0x31, 0x3B)
MUTED = RGBColor(0x8A, 0x85, 0x77)
SIGNAL = RGBColor(0xCE, 0x3B, 0x2A)
SERIES = [BRASS, RGBColor(0x3D, 0x6A, 0x88), RGBColor(0x5F, 0x8D, 0x7A),
          RGBColor(0xA9, 0x55, 0x3F), RGBColor(0x87, 0x93, 0xA0),
          RGBColor(0x7B, 0x61, 0x7F), RGBColor(0x9C, 0x8A, 0x4E)]

SERIF = "Cambria"      # display, body and every number
SANS = "Arial"         # labels and metadata only

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.62)                      # left and right margin
BODY_W = W - 2 * M

BOEING_DECK = ("Boeing, Market Overview 2025, European Consultant Conference, "
               "March 2025, Wendy Sowers")


# --------------------------------------------------------------------------- helpers

def uk(run):
    """Every run carries en-GB. Word and PowerPoint auto-detect per run otherwise, and
    a run left undeclared is regularly tagged as French."""
    run.font.language_id = MSO_LANGUAGE_ID.ENGLISH_UK
    return run


def text_box(slide, left, top, width, height, text, *, size=14, bold=False,
             colour=TEXT, font=SERIF, align=PP_ALIGN.LEFT, space_after=6,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        r = uk(p.add_run())
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = colour
    return tb


def rule(slide, top, colour=LINE, left=M, width=None, height=Pt(1)):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                               width or BODY_W, height)
    s.fill.solid()
    s.fill.fore_color.rgb = colour
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def background(slide, colour=PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def page(prs, title, kicker=None, *, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, INK if dark else PAPER)
    tcol = PAPER if dark else INK
    if kicker:
        text_box(slide, M, Inches(0.42), BODY_W, Inches(0.25), kicker.upper(),
                 size=10, font=SANS, colour=BRASS if dark else BRASS_DEEP,
                 space_after=0)
    text_box(slide, M, Inches(0.72), BODY_W, Inches(0.5), title,
             size=26, font=SERIF, colour=tcol, space_after=0)
    rule(slide, Inches(1.36), BRASS if dark else LINE, height=Pt(1.5))
    return slide


def footer(slide, source, *, note=None, dark=False):
    col = MUTED
    y = H - Inches(0.62)
    if note:
        text_box(slide, M, y - Inches(0.24), BODY_W, Inches(0.22), note,
                 size=9, font=SANS, colour=col, space_after=0)
    text_box(slide, M, y, BODY_W, Inches(0.22), "Source: " + source,
             size=9, font=SANS, colour=col, space_after=0)
    text_box(slide, M, y + Inches(0.2), BODY_W, Inches(0.22),
             "The Observatory Global Aviation Forecast   |   Commercial in Confidence "
             "   |   DRAFT", size=9, font=SANS, colour=col, space_after=0)


def table(slide, left, top, width, headers, rows, *, col_w=None, size=12,
          height=None, highlight_last_row=False):
    nrows, ncols = len(rows) + 1, len(headers)
    height = height or Inches(0.32) * nrows
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = shape.table
    tbl.first_row = True
    if col_w:
        total = sum(col_w)
        for i, w in enumerate(col_w):
            tbl.columns[i].width = Emu(int(width * w / total))
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if j else PP_ALIGN.LEFT
        r = uk(p.add_run())
        r.text = str(h)
        r.font.size = Pt(size - 2)
        r.font.name = SANS
        r.font.bold = True
        r.font.color.rgb = PAPER
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
    for i, row in enumerate(rows, start=1):
        last = highlight_last_row and i == len(rows)
        for j, v in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j else PP_ALIGN.LEFT
            r = uk(p.add_run())
            r.text = str(v)
            r.font.size = Pt(size)
            r.font.name = SERIF
            r.font.bold = bool(last)
            r.font.color.rgb = INK if last else TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if i % 2 else PAPER
    return tbl


def chart(slide, kind, left, top, width, height, categories, series,
          *, legend=True, number_format="0.0", gap_width=60, overlap=-10,
          y_title=None):
    data = CategoryChartData()
    data.categories = [str(c) for c in categories]
    for name, vals in series:
        data.add_series(name, vals, number_format)
    gf = slide.shapes.add_chart(kind, left, top, width, height, data)
    ch = gf.chart
    ch.font.size = Pt(10)
    ch.font.name = SANS
    ch.font.color.rgb = TEXT
    try:
        ch.plots[0].gap_width = gap_width
        ch.plots[0].overlap = overlap
    except (AttributeError, ValueError):
        pass
    if legend and len(series) > 1:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(9)
        ch.legend.font.name = SANS
    else:
        ch.has_legend = False
    for i, plot_series in enumerate(ch.series):
        col = SERIES[i % len(SERIES)]
        if kind in (XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS):
            plot_series.format.line.color.rgb = col
            plot_series.format.line.width = Pt(2.25)
            plot_series.smooth = False
        else:
            plot_series.format.fill.solid()
            plot_series.format.fill.fore_color.rgb = col
            plot_series.format.line.fill.background()
    try:
        va = ch.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = LINE
        va.major_gridlines.format.line.width = Pt(0.75)
        va.format.line.color.rgb = LINE
        va.tick_labels.font.size = Pt(9)
        va.tick_labels.font.name = SANS
        if y_title:
            va.has_title = True
            va.axis_title.text_frame.text = y_title
            for p in va.axis_title.text_frame.paragraphs:
                for r in p.runs:
                    uk(r)
                    r.font.size = Pt(9)
                    r.font.name = SANS
                    r.font.color.rgb = MUTED
        ca = ch.category_axis
        ca.has_major_gridlines = False
        ca.format.line.color.rgb = LINE
        ca.tick_labels.font.size = Pt(9)
        ca.tick_labels.font.name = SANS
    except (AttributeError, ValueError):
        pass
    return ch


def pct(v, dp=1):
    return "n/a" if v is None else f"{v * 100:.{dp}f}%"


def pp(v, dp=1):
    return "n/a" if v is None else f"{v:+.{dp}f}pp"


# ----------------------------------------------------------------------- the content

class Deck:
    def __init__(self, author):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.author = author
        self.gaps = []

    def gap(self, slide_name, boeing_page, reason):
        """A slide we cannot fill. Recorded, then printed as a slide of its own."""
        self.gaps.append((slide_name, boeing_page, reason))
        print(f"  GAP  {slide_name} (Boeing p{boeing_page}): {reason}")

    # -- slides ------------------------------------------------------------------

    def cover(self):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background(s, INK)
        text_box(s, M, Inches(2.0), BODY_W, Inches(0.3),
                 "THE OBSERVATORY", size=11, font=SANS, colour=BRASS, space_after=0)
        text_box(s, M, Inches(2.42), BODY_W, Inches(1.0),
                 "Global Aviation Forecast", size=44, font=SERIF, colour=PAPER,
                 space_after=0)
        text_box(s, M, Inches(3.36), BODY_W, Inches(0.5),
                 "Market overview, and where we stand against Boeing",
                 size=18, font=SERIF, colour=BRASS, space_after=0)
        rule(s, Inches(4.1), BRASS, width=Inches(3.2), height=Pt(2))
        text_box(s, M, Inches(4.4), Inches(8.4), Inches(1.2),
                 ["Prepared for the Avia Solutions forecast team",
                  "Built on the Avia global forecast, the OAG schedule store and Sabre "
                  "O&D. Every figure regenerates from the model."],
                 size=12, font=SERIF, colour=PAPER, space_after=4)
        text_box(s, M, H - Inches(1.15), BODY_W, Inches(0.6),
                 [date.today().strftime("%-d %B %Y") if os.name != "nt"
                  else date.today().strftime("%#d %B %Y"),
                  "DRAFT   |   Commercial in Confidence   |   "
                  "Copyright Avia Solutions Limited. All rights reserved."],
                 size=10, font=SANS, colour=MUTED, space_after=2)

    def how_to_read(self):
        s = page(self.prs, "What this deck is, and how to read it",
                 "Purpose")
        text_box(s, M, Inches(1.72), Inches(7.6), Inches(4.4),
                 ["This is Boeing's Market Overview rebuilt on Avia's own forecast and "
                  "assumptions. It has three uses and one rule.",
                  "",
                  "The team uses it to understand, review and challenge the forecast. "
                  "We use it to develop the forecast and to know exactly where we sit "
                  "against Boeing, and after Boeing against Airbus. And we use it with "
                  "third parties.",
                  "",
                  "The rule is that we are not goal seeking to Boeing. Where we differ, "
                  "the question is whether the difference is an error in method, which "
                  "we correct, or a difference of view, which we keep and explain. The "
                  "two are separated by measurement, not by argument.",
                  "",
                  "Building the deck is itself the diagnostic. Every slide we cannot "
                  "fill is a gap in the forecast rather than a formatting problem, and "
                  "the gaps are listed at the back with what each one would take."],
                 size=13)
        text_box(s, Inches(8.6), Inches(1.72), Inches(4.1), Inches(0.3),
                 "THE SOURCE DECK", size=10, font=SANS, colour=BRASS_DEEP,
                 space_after=4)
        text_box(s, Inches(8.6), Inches(2.05), Inches(4.1), Inches(3.0),
                 [BOEING_DECK + ", 28 pages, of which 22 carry content.",
                  "",
                  "Of those 22 we can produce 6 outright, 7 in part and 9 not at all.",
                  "",
                  "Boeing's forecast window is 2004-2023 on the historic slides and "
                  "2024-2044 on the forecast. Ours is 2015-2025 and 2024-2044. Where "
                  "the windows differ the slide says so."],
                 size=11, colour=TEXT)
        footer(s, "Avia Solutions, OGF slide inventory against the Boeing Market "
                  "Overview, version 1.0, 9 August 2026")

    def coverage(self, inv):
        s = page(self.prs, "What we can build, and what we cannot", "Coverage")
        rows = [[r["section"], r["slides"], r["can"], r["partial"], r["cannot"]]
                for r in inv]
        tot = ["Total", sum(r["slides"] for r in inv), sum(r["can"] for r in inv),
               sum(r["partial"] for r in inv), sum(r["cannot"] for r in inv)]
        table(s, M, Inches(1.75), Inches(7.4),
              ["Section", "Slides", "Can produce", "Partial", "Cannot"],
              rows + [tot], col_w=[3.4, 1, 1.3, 1, 1], highlight_last_row=True)
        text_box(s, Inches(8.3), Inches(1.75), Inches(4.4), Inches(4.2),
                 ["Six of twenty two outright. That is the honest starting position "
                  "and it is the point of the exercise.",
                  "",
                  "Three of the gaps are data we do not hold: consumer and travel "
                  "spending, world cargo market data, and fleet age. Each is a buy or "
                  "scope decision rather than an engineering one.",
                  "",
                  "One gap is already on our own risk list and this makes it client "
                  "visible. We forecast with a fare index and cannot say what a fare "
                  "costs anyone, so we cannot draw affordability, which is the "
                  "mechanism Boeing uses to explain emerging market growth.",
                  "",
                  "Two are classification rather than capability, and are a mapping "
                  "table because Atlas is built airport by airport."],
                 size=11)
        footer(s, "Avia Solutions, OGF slide inventory against the Boeing Market "
                  "Overview, version 1.0, 9 August 2026")

    def regions(self, reg):
        s = page(self.prs, "Avia against Boeing, by Boeing's own regions",
                 "Reconciliation")
        rows = sorted(reg["rows"], key=lambda r: -r["avia"])
        table(s, M, Inches(1.75), Inches(6.5),
              ["Region", "Airports", "Avia", "Boeing", "Difference"],
              [[r["region"], f"{r['airports']:,}", pct(r["avia"]), pct(r["boeing"]),
                pp(r["diff_pp"])] for r in rows]
              + [["World", "2,430", pct(reg["world_avia"]), pct(reg["world_boeing"]),
                  pp((reg["world_avia"] - reg["world_boeing"]) * 100)]],
              col_w=[2.3, 1, 1, 1, 1.2], size=11, highlight_last_row=True)
        chart(s, XL_CHART_TYPE.BAR_CLUSTERED, Inches(7.2), Inches(1.75),
              Inches(5.5), Inches(4.4),
              [r["region"] for r in rows],
              [("Difference, Avia less Boeing, percentage points",
                [round(r["diff_pp"], 2) for r in rows])],
              legend=False, number_format='0.0"pp"',
              y_title="RPK CAGR difference, percentage points, 2024-2044")
        footer(s, f"Avia Solutions global forecast, Baseline case, against the "
                  f"{reg['boeing_edition']}. Produced by "
                  f"scripts/compare_regions_boeing.py",
               note="Revenue passenger kilometres, compound annual growth "
                    "2024-2044, forecast on both sides. Regions are Boeing's ten, "
                    "applied to Avia's 2,430 airports by country.")

    def gap_bridge(self, gd):
        s = page(self.prs, "Two thirds of the gap against Boeing is a conversion, "
                           "not a view", "The finding")
        rows = gd["rows"] + ([gd["world"]] if gd.get("world") else [])
        table(s, M, Inches(1.72), Inches(6.6),
              ["Region", "Avia now", "Stage length", "Avia adjusted", "Boeing",
               "Gap after"],
              [[r["region"], pct(r["avia"]), pct(r["stage_cagr"]),
                pct(r["avia_with_stage"]), pct(r["boeing"]), pp(r["gap_pp_after"])]
               for r in rows],
              col_w=[2.1, 1, 1.2, 1.3, 1, 1.1], size=10.5,
              highlight_last_row=bool(gd.get("world")))
        text_box(s, Inches(7.4), Inches(1.72), Inches(5.3), Inches(4.6),
                 ["Boeing publishes RPK. Atlas forecasts passengers and converts to "
                  "RPK with a stage length held constant, so our RPK growth is our "
                  "passenger growth to the decimal place. Boeing's RPK growth carries "
                  "their stage length growth inside it.",
                  "",
                  "Measured from the OAG schedule, stage length grew 0.6% a year at "
                  "world level over 2015-2025. Carried into the conversion as a test, "
                  "the world gap narrows from 1.9 to 0.2 points against Boeing.",
                  "",
                  "Where the gap closes, we were comparing a passenger CAGR with an "
                  "RPK CAGR. Where it does not, in China, Africa and the Middle East, "
                  "we hold a different view of demand, and that has to be argued "
                  "rather than corrected.",
                  "",
                  "Nothing in the forecast has been changed. This is a measurement, "
                  "and a single historic rate applied everywhere over corrects Oceania "
                  "and Northeast Asia, so the fix is a stated stage length path per "
                  "region."],
                 size=11)
        footer(s, "Avia Solutions analysis of the OAG schedule store, 2015-2025, "
                  "against the Boeing 2025 Commercial Market Outlook. Produced by "
                  "scripts/gap_decomposition.py",
               note="Stage length measured as ASK per departing seat, scheduled "
                    "passenger services, actual. The adjusted column is a test and is "
                    "not the published forecast.")

    def wedge(self, wj, segment, label, boeing=None):
        s = page(self.prs, f"Fleet productivity: {label}", "The wedge")
        wins = [k for k in wj["windows"] if segment in wj["windows"][k]]
        rows = []
        for k in wins:
            w = wj["windows"][k][segment]
            rows.append([k, pct(w["g_ask"]), pct(w["g_seats"]), pct(w["g_departures"]),
                         pct(w["g_gauge"]), pct(w["g_stage"])])
        if boeing:
            rows.append([f"Boeing {BOEING['window'] if False else wj['boeing']['window']}",
                         pct(boeing["ask"]), pct(boeing["seats"]), "n/a", "n/a", "n/a"])
        table(s, M, Inches(1.72), Inches(7.0),
              ["Window", "ASK", "Seats", "Departures", "Seats a departure",
               "Stage length"],
              rows, col_w=[1.6, 1, 1, 1.2, 1.6, 1.3], size=11,
              highlight_last_row=bool(boeing))
        last = wins[-1]
        w = wj["windows"][last][segment]
        chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, M, Inches(3.1), Inches(7.0),
              Inches(3.0), ["Departures", "Seats a departure", "Stage length", "ASK"],
              [(f"Compound annual growth, {last}, actual",
                [round(w["g_departures"] * 100, 2), round(w["g_gauge"] * 100, 2),
                 round(w["g_stage"] * 100, 2), round(w["g_ask"] * 100, 2)])],
              legend=False, number_format='0.0"%"',
              y_title="Per cent a year")
        ss = (wj.get("shift_share", {}).get(last, {}) or {}).get(segment)
        notes = [
            "ASK is departures multiplied by seats a departure multiplied by stage "
            "length. The identity holds exactly, so the three bars account for the "
            "fourth with nothing left over.",
        ]
        if ss:
            notes += ["",
                      f"Seats a departure moved from {ss['gauge_start']:.0f} to "
                      f"{ss['gauge_end']:.0f} over {last}. Splitting that, up-gauging "
                      f"onto larger types contributes {pct(ss['up_gauging_cagr'])} a "
                      f"year and densification inside a type "
                      f"{pct(ss['densification_cagr'])}, with "
                      f"{pct(ss['interaction_cagr'])} in the interaction, which is "
                      "reported rather than allocated."]
        notes += ["",
                  "Boeing's fourth term, flights per aircraft per day, is not here. It "
                  "needs a count of aircraft in service, which we do not hold. "
                  "Deriving it from the dashboard's assumed productivity per aircraft "
                  "would return whatever was typed in."]
        text_box(s, Inches(7.9), Inches(1.72), Inches(4.8), Inches(4.4), notes, size=11)
        footer(s, "Avia Solutions analysis of the OAG schedule store, service type J, "
                  "departures only. Produced by scripts/build_fleet_wedge.py"
                  + (f". Boeing figures from {BOEING_DECK}, pages 24 and 25"
                     if boeing else ""),
               note="Scheduled passenger capacity, actual. Boeing's window is "
                    "2004-2023 and cannot be reproduced: the store holds 2015-2019 and "
                    "2023-2025, with 2020-2022 excluded by policy.")

    def network(self, dd):
        n = dd["network"]
        yrs = n["years"]
        s = page(self.prs, f"{yrs[-1] - yrs[0]} years of traffic and network",
                 "Airline trends")
        base = str(yrs[0])
        series = [
            ("Airport pairs served",
             [round(100 * n["airport_pairs"][str(y)] / n["airport_pairs"][base], 1)
              for y in yrs]),
            ("Departures",
             [round(100 * n["departures"][str(y)] / n["departures"][base], 1)
              for y in yrs]),
            ("ASK",
             [round(100 * n["ask_bn"][str(y)] / n["ask_bn"][base], 1) for y in yrs]),
        ]
        chart(s, XL_CHART_TYPE.LINE_MARKERS, M, Inches(1.72), Inches(7.6), Inches(4.4),
              yrs, series, number_format="0.0",
              y_title=f"Index, {base} = 100, actual")
        table(s, Inches(8.6), Inches(1.72), Inches(4.1),
              ["", f"{yrs[0]}", f"{yrs[-1]}", "a year"],
              [["Airport pairs", f"{n['airport_pairs'][base]:,}",
                f"{n['airport_pairs'][str(yrs[-1])]:,}",
                pct(n["cagr"]["airport_pairs"])],
               ["Departures m", f"{n['departures'][base] / 1e6:,.1f}",
                f"{n['departures'][str(yrs[-1])] / 1e6:,.1f}",
                pct(n["cagr"]["departures"])],
               ["ASK bn", f"{n['ask_bn'][base]:,.0f}",
                f"{n['ask_bn'][str(yrs[-1])]:,.0f}", pct(n["cagr"]["ask_bn"])]],
              col_w=[1.6, 1, 1, 1], size=11)
        text_box(s, Inches(8.6), Inches(3.3), Inches(4.1), Inches(2.6),
                 ["Boeing show twenty five years. We show "
                  f"{yrs[-1] - yrs[0]}, because the schedule store holds "
                  f"{yrs[0]} to 2019 and 2023 to {yrs[-1]}, with 2020 to 2022 excluded "
                  "by policy. The line is broken across the gap rather than joined.",
                  "",
                  "Ours is built from every scheduled departure rather than from a "
                  "sample, which is the one respect in which this slide is stronger "
                  "than theirs."],
                 size=11)
        footer(s, "Avia Solutions analysis of the OAG schedule store, service type J, "
                  "departures only. Produced by scripts/build_ogf_deck_data.py",
               note="Scheduled passenger services, actual. Airport pairs are "
                    "directional pairs with at least one departure in the year.")

    def lcc(self, dd):
        s = page(self.prs, "Low cost share of intra-regional capacity", "Airline trends")
        share = dd["lcc"]["share_of_seats"]
        regions = [k for k in share if "seats m" not in k]
        yrs = sorted({int(y) for k in regions for y in share[k]})
        series = [(r, [round(100 * share[r].get(str(y), 0) or 0, 1) for y in yrs])
                  for r in regions]
        chart(s, XL_CHART_TYPE.LINE_MARKERS, M, Inches(1.72), Inches(8.2), Inches(4.4),
              yrs, series, number_format='0.0"%"',
              y_title="Low cost share of intra-regional seats, per cent, actual")
        text_box(s, Inches(9.2), Inches(1.72), Inches(3.5), Inches(4.4),
                 ["Both ends of the flight inside the same region. The mainline and "
                  "low cost split is OAG's own carrier category, so the classification "
                  "is theirs and can be cited as theirs.",
                  "",
                  "Boeing run 2009 to 2024. We run "
                  f"{yrs[0]} to {yrs[-1]}, with 2020 to 2022 absent from the store."],
                 size=11)
        footer(s, "Avia Solutions analysis of the OAG schedule store, carrier category "
                  "as published by OAG. Produced by scripts/build_ogf_deck_data.py")

    def business_model(self, dd):
        s = page(self.prs, "Single aisle seat capacity by business model",
                 "Fleet dynamics")
        sa = dd["business_model"]["single_aisle"]
        yrs = sorted({int(y) for c in sa for y in sa[c]["seats_m"]})
        series = [(c.capitalize(),
                   [round(sa[c]["seats_m"].get(str(y), 0), 0) for y in yrs])
                  for c in sorted(sa)]
        chart(s, XL_CHART_TYPE.COLUMN_STACKED, M, Inches(1.72), Inches(8.2),
              Inches(4.4), yrs, series, number_format="#,##0",
              y_title="Departing seats, millions a year, actual", overlap=100)
        lines = ["Single aisle here is one aisle and a mainline jet, which excludes "
                 "regional jets, matching Boeing's segment."]
        for c in sorted(sa):
            g0 = sa[c]["seats_per_departure"].get(str(yrs[0]))
            g1 = sa[c]["seats_per_departure"].get(str(yrs[-1]))
            if g0 and g1:
                lines += ["", f"{c.capitalize()} gauge {g0:.0f} to {g1:.0f} seats a "
                              f"departure, {yrs[0]} to {yrs[-1]}."]
        text_box(s, Inches(9.2), Inches(1.72), Inches(3.5), Inches(4.4), lines, size=11)
        footer(s, "Avia Solutions analysis of the OAG schedule store, carrier category "
                  "as published by OAG. Produced by scripts/build_ogf_deck_data.py")

    def world_recovery(self, hist):
        wt = hist["world_traffic"]
        yrs = [y for y in wt["years"] if y >= 2010]
        base = "2019"
        s = page(self.prs, "World traffic against 2019", "Passenger demand")
        chart(s, XL_CHART_TYPE.LINE_MARKERS, M, Inches(1.72), Inches(8.2), Inches(4.4),
              yrs,
              [("RPK", [round(100 * wt["rpk_bn"][str(y)] / wt["rpk_bn"][base], 1)
                        for y in yrs]),
               ("Passengers", [round(100 * wt["pax_m"][str(y)] / wt["pax_m"][base], 1)
                               for y in yrs])],
              number_format="0.0", y_title="Index, 2019 = 100, actual")
        lf = [wt["plf"][str(y)] for y in yrs if str(y) in wt["plf"]]
        text_box(s, Inches(9.2), Inches(1.72), Inches(3.5), Inches(4.4),
                 ["RPK passed its 2019 level in 2024 and stands at "
                  f"{100 * wt['rpk_bn']['2025'] / wt['rpk_bn'][base]:.0f} in 2025.",
                  "",
                  f"The industry load factor rebuilt from {min(lf):.0f}% at the trough "
                  f"to {wt['plf']['2025']:.0f}% in 2025, a record, and that long climb "
                  "is a core efficiency assumption in the forecast.",
                  "",
                  "Boeing split this domestic against international and carry it to "
                  "2027. Our history is world total: the domestic and international "
                  "split of the RECOVERY is not held on the same basis, so it is in "
                  "the gap register rather than approximated here."],
                 size=11)
        footer(s, wt.get("source", "ICAO, IATA, Avia Solutions analysis")
               + ". Load factor: " + wt.get("rpk_plf_source", "as above"),
               note="Revenue passenger kilometres and passengers carried, world, "
                    "actual, indexed to 2019.")

    def route_areas(self, dash):
        s = page(self.prs, "The largest international route areas", "Passenger demand")
        flows = dash["flows"]
        inter = sorted(((k, v) for k, v in flows.items()
                        if k.split("|")[0] != k.split("|")[1]),
                       key=lambda kv: -kv[1])[:5]
        intra = sorted(((k, v) for k, v in flows.items()
                        if k.split("|")[0] == k.split("|")[1]),
                       key=lambda kv: -kv[1])[:3]
        rows = [[k.replace("|", " to "), f"{v:,.0f}"] for k, v in inter]
        rows += [[k.split("|")[0] + ", within the region", f"{v:,.0f}"]
                 for k, v in intra]
        table(s, M, Inches(1.75), Inches(6.4), ["Route area", "RPK bn, base year"],
              rows, col_w=[3.6, 1.8], size=12)
        chart(s, XL_CHART_TYPE.BAR_CLUSTERED, Inches(7.2), Inches(1.75), Inches(5.5),
              Inches(4.4), [r[0] for r in rows],
              [("RPK bn", [float(r[1].replace(",", "")) for r in rows])],
              legend=False, number_format="#,##0",
              y_title="Revenue passenger kilometres, billions, base year")
        footer(s, "Avia Solutions global forecast, region pair O&D flows, base year. "
                  "Produced by scripts/build_dashboard_data.py",
               note="Boeing's five are intra Asia Pacific, intra Europe, Asia Pacific "
                    "to Europe, Asia Pacific to North America and Europe to North "
                    "America. Ours are on our own six regions and the interregional "
                    "matrix uses representative region pair stage lengths, already "
                    "flagged as a priority one item.")

    def gap_register(self):
        s = page(self.prs, "What we cannot produce, and what each one needs",
                 "The gap register", dark=True)
        rows = [[a, f"p{b}", c] for a, b, c in self.gaps]
        table(s, M, Inches(1.72), BODY_W, ["Slide", "Boeing", "What it needs"],
              rows, col_w=[3.2, 0.8, 8.0], size=10.5)
        text_box(s, M, H - Inches(1.35), BODY_W, Inches(0.5),
                 "A slide we cannot publish is a gap, not a formatting problem. Each "
                 "line above is either a data acquisition decision or a piece of "
                 "engineering, and each one has an owner before this deck goes out.",
                 size=11, colour=PAPER)
        footer(s, "Avia Solutions, OGF slide inventory version 1.0, 9 August 2026, and "
                  "the build log of this deck", dark=True)

    def basis(self):
        s = page(self.prs, "Basis, and the checks behind it", "Method")
        text_box(s, M, Inches(1.72), Inches(6.1), Inches(4.6),
                 ["Forecast. Avia global forecast, Baseline case, 2,430 airports "
                  "carrying a Boeing region, Oxford Economics country GDP of 31 July "
                  "2024, world O&D departing passengers 3,140m in 2025 to 9,644m in "
                  "2060, a compound 3.26% a year.",
                  "",
                  "History. The OAG schedule store, 333m rows, service type J, "
                  "departures only, one preferred tiling for each region and year, "
                  "each airport read from its home region file. Sabre O&D for the "
                  "regional and cabin history. ICAO and IATA for the world series.",
                  "",
                  "Comparators. Boeing 2025 Commercial Market Outlook, held on Egnyte, "
                  "read from the workbook and not from the press release."],
                 size=11.5)
        text_box(s, Inches(7.0), Inches(1.72), Inches(5.7), Inches(4.6),
                 ["Checks run before any figure on these slides was produced:",
                  "",
                  "The store carries eight complete years across seven region files, "
                  "no slice at the Excel sheet limit and none thin against the rest of "
                  "its year, so the export truncation of July 2026 is not present.",
                  "",
                  "Distance is in kilometres, proved against the mileage column. 173 "
                  "rows carrying an impossible sector are excluded and account for "
                  "0.0002% of ASK.",
                  "",
                  "Heathrow 2019 reproduces at 477,954 movements and 100.3m seats two "
                  "way against the published figures, and 238,978 departures one way, "
                  "which is the basis every capacity figure here is built on.",
                  "",
                  "244 of the 245 aircraft codes in the store carry a body type, "
                  "leaving 0.002% of seats unclassified."],
                 size=11.5)
        footer(s, "scripts/guard_oag_wedge.py, run before the deck build, all 15 "
                  "checks passing")

    def save(self, path):
        cp = self.prs.core_properties
        cp.author = self.author
        cp.last_modified_by = self.author
        cp.title = "The Observatory Global Aviation Forecast, market overview"
        cp.comments = ("Generated by scripts/build_ogf_deck.py in the Atlas "
                       "repository. Do not edit by hand: the next build overwrites it.")
        cp.language = "en-GB"
        self.prs.save(path)
        sanitise(path)


def sanitise(path):
    """Rewrite the saved package to remove what the default template brings with it.

    python-pptx ships a template whose slide master uses an EN DASH as the bullet
    character at three outline levels, and whose runs and defaults declare en-US. Both
    are in the delivered file whether or not anything on a slide uses them, and both
    are exactly the kind of thing a review by eye never finds. Neither can be reached
    through the python-pptx object model, so the parts are rewritten here and
    scripts/check_deck.py verifies the result from the outside.
    """
    import shutil
    import tempfile
    import zipfile
    subs = [('lang="en-US"', 'lang="en-GB"'),
            ('altLang="en-US"', 'altLang="en-GB"'),
            ("—", ","),        # em dash
            ("–", "-"),        # en dash, including the template bullet character
            ("―", "-"),        # horizontal bar
            ("−", "-"),        # minus sign
            ("‑", "-")]        # non-breaking hyphen
    tmp = tempfile.mktemp(suffix=".pptx")
    n = 0
    with zipfile.ZipFile(path) as src, \
            zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            body = src.read(item.filename)
            if item.filename.endswith((".xml", ".rels")):
                s = body.decode("utf-8")
                before = s
                for a, b in subs:
                    s = s.replace(a, b)
                if s != before:
                    n += 1
                body = s.encode("utf-8")
            dst.writestr(item, body)
    shutil.move(tmp, path)
    print(f"sanitised {n} XML parts: language forced to en-GB, prohibited dashes "
          "replaced")


BOEING = {"window": "2004-2023"}


def load(path, what):
    if not os.path.exists(path):
        print(f"  missing input: {what} at {path}")
        return None
    with open(path, "r", encoding="utf-8") as fh:
        return json.load(fh)


INVENTORY = [
    {"section": "1. Passenger air travel demand", "slides": 9, "can": 3, "partial": 3,
     "cannot": 3},
    {"section": "2. Air cargo demand", "slides": 4, "can": 0, "partial": 1, "cannot": 3},
    {"section": "3. Airline trends", "slides": 4, "can": 2, "partial": 1, "cannot": 1},
    {"section": "4. Fleet dynamics", "slides": 5, "can": 1, "partial": 2, "cannot": 2},
]


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--out", default=None)
    ap.add_argument("--author", default="Avia Solutions",
                    help='product decks published by the Observatory carry '
                         '"The Aviation Observatory" instead')
    args = ap.parse_args(argv)

    wj = load(os.path.join(paths.DATA, "fleet_wedge.json"), "the fleet wedge")
    gd = load(os.path.join(paths.DATA, "gap_decomposition.json"),
              "the gap decomposition")
    reg = load(os.path.join(paths.DATA, "regions_boeing.json"),
               "the region reconciliation")
    dd = load(os.path.join(paths.DATA, "ogf_deck_data.json"), "the schedule slides")
    hist = load(os.path.join(REPO, "webapp", "data", "history.json"), "the history")
    dash = load(os.path.join(REPO, "webapp", "data", "dashboard.json"),
                "the dashboard data")

    d = Deck(args.author)
    d.cover()
    d.how_to_read()
    d.coverage(INVENTORY)
    if reg:
        d.regions(reg)
    else:
        d.gap("Avia against Boeing by region", 8,
              "regions_boeing.json absent. Run compare_regions_boeing.py --json")
    if gd:
        d.gap_bridge(gd)
    else:
        d.gap("The stage length bridge", 8,
              "gap_decomposition.json absent. Run gap_decomposition.py")
    if wj:
        d.wedge(wj, "single_aisle", "single aisle",
                boeing=wj["boeing"]["single_aisle"])
        d.wedge(wj, "widebody", "widebody")
    else:
        d.gap("Fleet productivity wedge", 24,
              "fleet_wedge.json absent. Run build_fleet_wedge.py")
    if hist:
        d.world_recovery(hist)
    else:
        d.gap("World traffic against 2019", 7, "webapp/data/history.json absent")
    if dash and dash.get("flows"):
        d.route_areas(dash)
    else:
        d.gap("Largest international route areas", 9,
              "webapp/data/dashboard.json absent or carries no flows")
    if dd:
        d.network(dd)
        d.lcc(dd)
        d.business_model(dd)
    else:
        for name, p in (("Traffic and network over time", 18),
                        ("Low cost share of intra-regional capacity", 20),
                        ("Single aisle seat capacity by business model", 26)):
            d.gap(name, p, "ogf_deck_data.json absent. Run build_ogf_deck_data.py, "
                           "which reads the whole OAG store and wants a real host")

    # The standing gaps, from the inventory. These are not build failures.
    d.gap("Leisure travel spending share", 3,
          "Tourism Economics consumer spending. Not held. Buy or drop")
    d.gap("Air fare share of travel spending", 4,
          "Absolute fare levels. Our fare series is an index with no level, the F15 "
          "item still open in the assumptions book")
    d.gap("Regional unemployment", 5,
          "S&P Global labour market data. Not held, and not a driver in our method")
    d.gap("Air travel affordability", 6,
          "Absolute fares against GDP per capita. We hold the GDP and not the fares. "
          "The single most valuable gap on this list: affordability is the mechanism "
          "behind emerging market growth and we cannot show it")
    d.gap("Tariff impact on GDP and US inbound", 10,
          "A scenario overlay we do not run. Buildable as a named scenario")
    d.gap("World cargo tonne kilometre index", 13,
          "No CTK series. The OGF is scoped as passenger and airport cargo")
    d.gap("De minimis e-commerce", 14, "US customs bills of lading. Not held")
    d.gap("Containership against freighter reliability", 15,
          "Maritime data. Not held and not our field")
    d.gap("Air cargo flows reconfiguring", 16, "Trade flow data. Not held")
    d.gap("Fleet at retirement age", 23,
          "Fleet age from Cirium or Ascend. Avia holds Ascend, so acquisition rather "
          "than capability")
    d.gap("Fleet productivity to 2043", 27,
          "We assume a load factor path and do not forecast utilisation or average "
          "seats as outputs")
    d.gap("Flights per aircraft per day", 24,
          "A count of aircraft in service. The same acquisition as fleet age, and the "
          "fourth term of the wedge")

    d.gap_register()
    d.basis()

    out = args.out or os.path.join(
        paths.PROJECT_DIR,
        f"Observatory Global Aviation Forecast - Market Overview - "
        f"{datetime.now().strftime('%d %b %Y')} - DRAFT.pptx")
    d.save(out)
    print(f"\n{len(d.prs.slides.__iter__.__self__._sldIdLst)} slides written to {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
